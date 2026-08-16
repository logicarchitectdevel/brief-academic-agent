import os
import time
import streamlit as st
from google import genai
from google.genai import types
from pypdf import PdfReader
from pptx import Presentation
from gtts import gTTS

# =====================================================================
# 🚨 PASTE YOUR API KEY DIRECTLY BETWEEN THE QUOTES BELOW 🚨
# =====================================================================
api_key = "REAL_API_KEY_HERE"

# Setup Streamlit Page Configuration with Robot Icon
st.set_page_config(
    page_title="Brief | Academic Helping Agent", 
    page_icon="🤖", 
    layout="wide",
    initial_sidebar_state="expanded"
)

# 🎨 PREMIUM UI: Seamless color-matched typography, exact vertical centering, and clean states
st.markdown("""
    <style>
        /* 1. Alignment Container for the Sidebar Branding */
        .brand-header-container {
            display: flex !important;
            flex-direction: column !important;
            align-items: flex-start !important;
            gap: 0px !important;
            margin-top: -15px !important;
            margin-bottom: 0px !important;
        }
        .brand-row {
            display: flex !important;
            align-items: center !important;
            gap: 14px !important;
        }
        /* SOLUTION 1 DESIGN: Safely forces the native emoji to be a richer, darker royal blue (#1e3a8a) */
        .sidebar-robot {
            font-size: 50px !important;
            line-height: 1 !important;
            margin-top: 6px !important;
            display: inline-block !important;
            color: transparent !important;
            text-shadow: 0 0 0 #1e3a8a !important;
        }
        /* 2. Custom Typography with Gradient */
        .premium-brand-title {
            font-family: 'Inter', 'Google Sans', 'Segoe UI', sans-serif !important;
            font-size: 34px !important;
            font-weight: 800 !important;
            background: linear-gradient(135deg, #2563eb 0%, #1e3a8a 100%) !important;
            -webkit-background-clip: text !important;
            -webkit-text-fill-color: transparent !important;
            letter-spacing: -0.5px !important;
            display: inline-block !important;
            margin: 0 !important;
            line-height: 1.1 !important;
        }
        .premium-brand-subtitle {
            font-family: 'Inter', sans-serif !important;
            font-size: 13px !important;
            font-weight: 500 !important;
            color: #6b7280 !important;
            margin-top: 6px !important;
            margin-bottom: 20px !important;
            padding-left: 2px !important;
        }
        /* 3. Primary Button - Lighter blue base, Darker blue hover */
        div.stButton > button[kind="primary"] {
            background-color: #2563eb !important; 
            border-color: #2563eb !important;
            color: white !important;
            font-weight: 600 !important;
            transition: all 0.25s ease-in-out !important;
        }
        div.stButton > button[kind="primary"]:hover {
            background-color: #1e3a8a !important; 
            border-color: #1e3a8a !important;
            box-shadow: 0 4px 12px rgba(30, 58, 138, 0.2) !important;
        }
        /* Hides the annoying hover link icons next to headers */
        .viewerBadge_container__1QS1Z, a.header-anchor {
            display: none !important;
        }
        h1 a, h2 a, h3 a, h4 a, h5 a, h6 a {
            display: none !important;
        }
    </style>
""", unsafe_allow_html=True)

# Safe API Client Initialization
client = None
if api_key and "YOUR_REAL_KEY" not in api_key:
    client = genai.Client(api_key=api_key)

def extract_text_from_pdf(pdf_file):
    try:
        reader = PdfReader(pdf_file)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text
    except Exception as e:
        st.error(f"Failed to read PDF file: {e}")
        return ""

def extract_text_from_pptx(pdf_file):
    try:
        prs = Presentation(pdf_file)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Failed to read PPTX file: {e}")
        return ""
def run_brief_agent(course_material: str) -> str:
    if not client:
        return "❌ API Client not configured. Please check your API Key."
    
    system_instruction = (
        "You are Brief, an official autonomous academic helping agent. "
        "Your role is to analyze university course materials and extract direct knowledge. "
        "CRITICAL FORMAT RULES FOR VISUAL STRUCTURE: "
        "Provide an extensive, comprehensive report that balances precise crispness with deep academic substance. "
        "The response must be heavily structured VISUALLY so it is extremely easy to scan. "
        "Use markdown elements aggressively: clear thematic section headers, short scannable paragraphs, distinct bulleted fragments, and bold key terms. "
        "Every sentence must be packed with facts, direct definitions, theoretical context, or key data, but separated layout-wise for readability. "
        "Structure: "
        "1. Executive Summary: A profound, highly explicit, extensive explanation of the core definitions and concepts, highly broken down into logical sections or clear subheadings. "
        "2. Action Plan: A comprehensive, high-value sequential To-Do list with explicit guidelines for the student. "
        "Do NOT use code block boxes, raw markdown formatting hacks, or weird text highlighting placeholders."
    )
    
    try:
        response = client.models.generate_content(
            model='gemini-3.5-flash',
            contents=f"Analyze this material thoroughly. Provide a highly explicit, extensive academic structural synthesis with exceptional visual structure: {course_material}",
            config=types.GenerateContentConfig(
                system_instruction=system_instruction,
                temperature=0.2,
                max_output_tokens=8192,
                top_p=0.95
            )
        )
        return response.text
    except Exception as e:
        return f"❌ AI Generation Error: {str(e)}"

def convert_text_to_audio(text_content, voice_tld):
    if "❌ AI Generation Error" in text_content:
        return None
    try:
        clean_text = text_content.replace("#", "").replace("*", "").replace("- [ ]", "Task:")
        tts = gTTS(text=clean_text, lang='en', tld=voice_tld, slow=False)
        temp_filename = "temp_brief_podcast.mp3"
        tts.save(temp_filename)
        with open(temp_filename, "rb") as f:
            audio_bytes = f.read()
        os.remove(temp_filename)
        return audio_bytes
    except Exception as e:
        st.error(f"Audio Generation Error: {e}")
        return None

# =====================================================================
# 🌐 PREMIUM STREAMLIT INTERFACE (UI)
# =====================================================================
st.sidebar.markdown("""
    <div class="brand-header-container">
        <div class="brand-row">
            <div class="sidebar-robot">🤖</div>
            <h1 class="premium-brand-title">Brief</h1>
        </div>
        <div class="premium-brand-subtitle">Academic Helping Agent</div>
    </div>
""", unsafe_allow_html=True)

st.sidebar.markdown("---")

st.sidebar.markdown("#### 🎙️ Voice Narrator Accent")
selected_accent = st.sidebar.selectbox(
    "Choose your preferred voice:",
    ["British Accent (Premium Academic)", "American Accent (Standard Modern)"]
)

tld_mapping = {
    "British Accent (Premium Academic)": "co.uk",
    "American Accent (Standard Modern)": "com"
}
chosen_tld = tld_mapping[selected_accent]

st.sidebar.markdown("---")

# Prestigous minimalist title
st.title("Your Academic Workspace")
st.markdown("##### *Instantly transform heavy university lectures into structured reports and premium audio podcasts.*")
st.divider()

if not client:
    st.warning("🚨 Please insert your real Gemini API key inside the script to activate the AI Agent.")
else:
    uploaded_file = st.file_uploader("📂 Upload your document here (Supports: TXT, PDF, PPTX)", type=["txt", "pdf", "pptx"])

    if uploaded_file is not None:
        file_name = uploaded_file.name
        name_part, ext_part = os.path.splitext(file_name)
        
        st.write("") 
        
        if st.button("🚀 Execute Autonomous Analysis", type="primary", use_container_width=True):
            with st.spinner("🤖 Agent is executing background reading, synthesis, and podcast compilation..."):
                
                file_ext = ext_part.lower()
                extracted_content = ""
                
                if file_ext == ".txt":
                    extracted_content = uploaded_file.read().decode("utf-8")
                elif file_ext == ".pdf":
                    extracted_content = extract_text_from_pdf(uploaded_file)
                elif file_ext == ".pptx":
                    extracted_content = extract_text_from_pptx(uploaded_file)

                if extracted_content.strip():
                    ai_analysis_report = run_brief_agent(extracted_content)
                    audio_podcast_data = convert_text_to_audio(ai_analysis_report, chosen_tld)
                    
                    st.toast("Processing Complete! ✨", icon="✅")
                    st.divider()
                    
                    st.subheader("🎙️ Audio Brief Podcast")
                    with st.container(border=True):
                        if audio_podcast_data:
                            st.audio(audio_podcast_data, format="audio/mp3")
                            st.write("")
                            st.download_button(
                                label="💾 Download Podcast (MP3 Format)",
                                data=audio_podcast_data,
                                file_name=f"{name_part}_podcast.mp3",
                                mime="audio/mp3",
                                use_container_width=True
                            )
                    
                    st.write("") 
                    st.divider()
                    
                    st.subheader("Structured Academic Report")
                    with st.container(border=True):
                        st.markdown(ai_analysis_report)
                else:
                    st.error("❌ The uploaded file appears to be empty or unreadable.")
